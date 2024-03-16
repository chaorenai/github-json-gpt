from pptx import Presentation
import json

def ppt_to_json(ppt_file, json_file):
    prs = Presentation(ppt_file)
    data = {'content': [slide.notes_slide.notes_text_frame.text for slide in prs.slides]}
    with open(json_file, 'w') as jf:
        json.dump(data, jf, indent=4)

ppt_to_json('path/to/your/example.pptx', 'path/to/your/example.json')